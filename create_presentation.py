#!/usr/bin/env python3
"""
THE TRAITORS - Game Night Rules Presentation Generator
=====================================================
Generates a dramatic, castle-themed PowerPoint presentation for a 34-player
Traitors game night. Uses dark moody aesthetics inspired by the TV show.

Usage:
    python3 create_presentation.py

Output:
    traitors_rules.pptx in the same directory

Requires:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==============================================================================
# THEME CONSTANTS
# ==============================================================================

# Colors
BURGUNDY = RGBColor(0x8B, 0x00, 0x00)
DARK_BURGUNDY = RGBColor(0x5C, 0x00, 0x00)
GOLD = RGBColor(0xDA, 0xA5, 0x20)
PALE_GOLD = RGBColor(0xF0, 0xD0, 0x70)
DARK_GREEN = RGBColor(0x01, 0x32, 0x20)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
BLACK = RGBColor(0x0D, 0x0D, 0x0D)
NEAR_BLACK = RGBColor(0x12, 0x12, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF5, 0xF0, 0xE0)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x99, 0x99, 0x99)
BLOOD_RED = RGBColor(0xAA, 0x00, 0x00)
DARK_SLATE = RGBColor(0x20, 0x20, 0x20)
MEDIUM_GREEN = RGBColor(0x02, 0x4A, 0x30)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Font sizes
TITLE_SIZE = Pt(48)
SUBTITLE_SIZE = Pt(28)
HEADING_SIZE = Pt(40)
BODY_SIZE = Pt(28)
SMALL_BODY_SIZE = Pt(24)
DETAIL_SIZE = Pt(20)
MEGA_TITLE_SIZE = Pt(72)
LARGE_TITLE_SIZE = Pt(56)


# ==============================================================================
# HELPER FUNCTIONS
# ==============================================================================

def set_slide_background(slide, color):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_full_slide_rect(slide, color, transparency=0):
    """Add a rectangle covering the entire slide as a background layer."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_decorative_bar(slide, y_pos, color=GOLD, width=None, height=Inches(0.04),
                       x_pos=None):
    """Add a thin decorative horizontal bar."""
    if width is None:
        width = Inches(8)
    if x_pos is None:
        x_pos = (SLIDE_WIDTH - width) / 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x_pos, y_pos, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_side_accent(slide, side="left", color=BURGUNDY):
    """Add a vertical accent bar on the left or right side."""
    width = Inches(0.15)
    if side == "left":
        x = 0
    else:
        x = SLIDE_WIDTH - width
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, 0, width, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_top_bottom_bars(slide, color=DARK_BURGUNDY, height=Inches(0.08)):
    """Add thin accent bars at top and bottom of slide."""
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, height
    )
    top.fill.solid()
    top.fill.fore_color.rgb = color
    top.line.fill.background()
    bot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - height, SLIDE_WIDTH, height
    )
    bot.fill.solid()
    bot.fill.fore_color.rgb = color
    bot.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=BODY_SIZE,
                 font_color=CREAM, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(6)
    if line_spacing:
        p.line_spacing = line_spacing
    txBox.text_frame.auto_size = None
    return txBox


def add_multi_text_box(slide, left, top, width, height, lines,
                       anchor=MSO_ANCHOR.TOP):
    """
    Add a text box with multiple formatted lines.
    lines: list of dicts with keys: text, size, color, bold, alignment, font,
           spacing_after, line_spacing
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line_info.get("text", "")
        p.font.size = line_info.get("size", BODY_SIZE)
        p.font.color.rgb = line_info.get("color", CREAM)
        p.font.bold = line_info.get("bold", False)
        p.font.name = line_info.get("font", "Calibri")
        p.alignment = line_info.get("alignment", PP_ALIGN.LEFT)
        p.space_after = line_info.get("spacing_after", Pt(8))
        ls = line_info.get("line_spacing", None)
        if ls:
            p.line_spacing = ls

    txBox.text_frame.auto_size = None
    return txBox


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def diamond_bullet(text):
    """Return text with a diamond bullet prefix."""
    return f"\u25C6  {text}"


def add_centered_heading(slide, text, y_pos=Inches(0.5), color=GOLD,
                          size=HEADING_SIZE, width=Inches(11)):
    """Add a centered heading at the specified position."""
    x_pos = (SLIDE_WIDTH - width) / 2
    return add_text_box(
        slide, x_pos, y_pos, width, Inches(1),
        text, font_size=size, font_color=color, bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )


def standard_slide_setup(slide, title_text, bg_color=NEAR_BLACK):
    """Standard setup: background, side accents, top/bottom bars, title, divider."""
    set_slide_background(slide, bg_color)
    add_side_accent(slide, "left", BURGUNDY)
    add_side_accent(slide, "right", BURGUNDY)
    add_top_bottom_bars(slide, DARK_BURGUNDY)
    add_centered_heading(slide, title_text, y_pos=Inches(0.4))
    add_decorative_bar(slide, Inches(1.15), GOLD, width=Inches(6))


# ==============================================================================
# SLIDE BUILDERS
# ==============================================================================

def build_slide_01_title(prs):
    """Slide 1: Title slide - THE TRAITORS in large gold text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, BLACK)

    # Subtle dark burgundy vignette rectangle
    vignette = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(0.5), SLIDE_WIDTH - Inches(2), SLIDE_HEIGHT - Inches(1)
    )
    vignette.fill.solid()
    vignette.fill.fore_color.rgb = RGBColor(0x10, 0x08, 0x08)
    vignette.line.fill.background()

    # Inner frame with gold border
    inner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(1), SLIDE_WIDTH - Inches(4), SLIDE_HEIGHT - Inches(2)
    )
    inner.fill.solid()
    inner.fill.fore_color.rgb = RGBColor(0x14, 0x0A, 0x0A)
    inner.line.color.rgb = GOLD
    inner.line.width = Pt(2)

    # Top decorative bar
    add_decorative_bar(slide, Inches(2.0), GOLD, width=Inches(5), height=Inches(0.03))

    # Main title: THE TRAITORS
    add_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(10)) / 2,
        top=Inches(2.2),
        width=Inches(10),
        height=Inches(1.8),
        text="THE TRAITORS",
        font_size=MEGA_TITLE_SIZE,
        font_color=GOLD,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri"
    )

    # Bottom decorative bar
    add_decorative_bar(slide, Inches(4.0), GOLD, width=Inches(5), height=Inches(0.03))

    # Subtitle block
    add_multi_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(8)) / 2,
        top=Inches(4.3),
        width=Inches(8),
        height=Inches(1.5),
        lines=[
            {"text": "Game Night", "size": Pt(36), "color": CREAM,
             "bold": False, "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(12)},
            {"text": "34 Players  |  2 Houses  |  1 Winner",
             "size": Pt(22), "color": DIM_GRAY, "bold": False,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(12)},
            {"text": "Trust No One.",
             "size": Pt(28), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_top_bottom_bars(slide, DARK_BURGUNDY, Inches(0.06))

    add_speaker_notes(slide, (
        "WELCOME EVERYONE. Wait for the room to settle. Let the title sit on screen "
        "for a moment -- build the atmosphere. Dim the lights if possible.\n\n"
        "Say: 'Welcome... to The Traitors. Tonight, 34 of you enter this game. "
        "Not all of you can be trusted. Some of you have already been chosen -- "
        "chosen to deceive, to manipulate, and to betray. The question is... who?'\n\n"
        "Pause. Let it land. Then advance."
    ))


def build_slide_02_welcome(prs):
    """Slide 2: Welcome to the Castle - atmospheric intro."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NEAR_BLACK)

    # Dark green atmospheric overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x08, 0x14, 0x0D)
    overlay.line.fill.background()

    add_side_accent(slide, "left", DARK_GREEN)
    add_side_accent(slide, "right", DARK_GREEN)
    add_top_bottom_bars(slide, DARK_GREEN, Inches(0.06))

    add_centered_heading(slide, "Welcome to the Castle", y_pos=Inches(1.0),
                          color=GOLD, size=LARGE_TITLE_SIZE)
    add_decorative_bar(slide, Inches(1.85), GOLD, width=Inches(4), height=Inches(0.03))

    add_multi_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(9)) / 2,
        top=Inches(2.3),
        width=Inches(9),
        height=Inches(4.5),
        lines=[
            {"text": "Behind these walls, alliances will form and shatter.",
             "size": Pt(28), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(24), "line_spacing": 1.4},
            {"text": "Trust will be weaponised.",
             "size": Pt(28), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(24), "line_spacing": 1.4},
            {"text": "Secrets will be kept -- and exposed.",
             "size": Pt(28), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(36), "line_spacing": 1.4},
            {"text": "Among you walk Traitors.",
             "size": Pt(32), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(16)},
            {"text": "They look like you. They talk like you. They sit beside you.",
             "size": Pt(24), "color": DIM_GRAY, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(24)},
            {"text": "But they want you gone.",
             "size": Pt(28), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "Read this slide slowly and dramatically. Low voice.\n\n"
        "'Welcome to the Castle. Tonight you enter a world of deception. "
        "Behind these walls, alliances will form -- and shatter. Trust will be "
        "weaponised. And among you, right now, in this room... walk Traitors. "
        "They look like you. They talk like you. They sit beside you. "
        "But they want you gone.'\n\n"
        "Pause after 'gone.' Let the paranoia begin."
    ))


def build_slide_03_the_game(prs):
    """Slide 3: The Game - brief overview of faithful vs traitors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    standard_slide_setup(slide, "The Game")

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "34 players enter. Not all are who they seem.",
             "size": Pt(30), "color": CREAM, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(28)},
            {"text": diamond_bullet("28 of you are FAITHFUL"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(8)},
            {"text": "     Your goal: find and banish the Traitors.",
             "size": SMALL_BODY_SIZE, "color": LIGHT_GRAY, "spacing_after": Pt(20)},
            {"text": diamond_bullet("6 of you are TRAITORS"),
             "size": BODY_SIZE, "color": BURGUNDY, "bold": True,
             "spacing_after": Pt(8)},
            {"text": "     Your goal: eliminate the Faithful without being caught.",
             "size": SMALL_BODY_SIZE, "color": LIGHT_GRAY, "spacing_after": Pt(28)},
            {"text": "Through challenges, debates, votes, and secret night-time",
             "size": BODY_SIZE, "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(4)},
            {"text": "murders, the game narrows from 34 down to a Final Circle.",
             "size": BODY_SIZE, "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(28)},
            {"text": "If all Traitors are banished, the Faithful win.",
             "size": Pt(26), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "If Traitors survive to the end, they win.",
             "size": Pt(26), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "Explain the core concept clearly.\n\n"
        "'This is a game of social deduction. 34 of you are playing tonight. "
        "28 of you are Faithful -- your job is to figure out who the Traitors are "
        "and vote them out. But 6 of you... are Traitors. You know who you are. "
        "Your job is to blend in, deceive, and eliminate the Faithful one by one -- "
        "without getting caught.\n\n"
        "Through challenges, heated debates, dramatic votes, and secret night-time "
        "murders, we will whittle this group down from 34 to a Final Circle of 6. "
        "If all Traitors are gone by then, the Faithful win. If even one Traitor "
        "survives... the Traitors win.'"
    ))


def build_slide_04_your_role(prs):
    """Slide 4: Your Role - sealed envelopes, role assignment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    standard_slide_setup(slide, "Your Role")

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "Each of you will receive a sealed envelope.",
             "size": Pt(30), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(28)},
            {"text": diamond_bullet("Do NOT open it until instructed"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(14)},
            {"text": diamond_bullet("Open it privately -- do not show anyone"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(14)},
            {"text": diamond_bullet("Read your role card silently"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(14)},
            {"text": diamond_bullet("Seal it back in the envelope"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(14)},
            {"text": diamond_bullet("Return the envelope to the Host"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(28)},
            {"text": "Do not react. Do not gasp. Do not smile.",
             "size": Pt(30), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(14)},
            {"text": "Everyone is watching.",
             "size": Pt(26), "color": DIM_GRAY,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "This is a critical moment. You need complete control of the room.\n\n"
        "'In a moment, I will hand each of you a sealed envelope. Do NOT open it "
        "until I say so. When I give the word, open your envelope privately. Read "
        "the card inside. It will tell you if you are Faithful... or a Traitor. "
        "Do not react. Do not gasp. Do not look at anyone. Do not smile. "
        "Everyone in this room is watching everyone else right now. "
        "Read it, memorise it, seal it back up, and return it to me.\n\n"
        "Remember: from this moment on, the game has begun.'\n\n"
        "LOGISTICS: Hand out envelopes. Give them 30 seconds. Collect them all back. "
        "Traitor envelopes also contain their house assignment and a reminder that "
        "they will meet during the first Night phase."
    ))


def build_slide_05_faithful(prs):
    """Slide 5: The Faithful - mission and abilities."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NEAR_BLACK)

    # Subtle blue-tinted background for faithful
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x18)
    overlay.line.fill.background()

    add_side_accent(slide, "left", GOLD)
    add_side_accent(slide, "right", GOLD)
    add_top_bottom_bars(slide, RGBColor(0x8B, 0x75, 0x00), Inches(0.06))

    add_centered_heading(slide, "The Faithful", y_pos=Inches(0.4), color=GOLD,
                          size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(1.15), GOLD, width=Inches(4))

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "You are loyal. You are honest. You are hunted.",
             "size": Pt(30), "color": PALE_GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(28)},
            {"text": "YOUR MISSION", "size": Pt(26), "color": GOLD, "bold": True,
             "spacing_after": Pt(14)},
            {"text": diamond_bullet("Identify the Traitors among you"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Vote to banish them at the Round Table"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Survive the night"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(24)},
            {"text": "YOUR TOOLS", "size": Pt(26), "color": GOLD, "bold": True,
             "spacing_after": Pt(14)},
            {"text": diamond_bullet("Your vote at every banishment"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Your powers of observation and persuasion"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Challenge victories may grant special powers"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
        ]
    )

    add_speaker_notes(slide, (
        "'If your card said FAITHFUL, this is your mission. You are loyal. You are "
        "honest. And you are being hunted. The Traitors are in this room, pretending "
        "to be just like you. Your job is to figure out who they are. Watch "
        "behaviour. Listen to arguments. Notice who deflects, who accuses too "
        "quickly, who stays too quiet. Your main weapon is your vote at the "
        "Round Table. Use it wisely. And try to survive the night -- because the "
        "Traitors will be choosing someone to eliminate after dark.'"
    ))


def build_slide_06_traitors(prs):
    """Slide 6: The Traitors - vague and mysterious, no secrets revealed."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)

    # Deep red atmospheric overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x18, 0x06, 0x06)
    overlay.line.fill.background()

    add_side_accent(slide, "left", BLOOD_RED)
    add_side_accent(slide, "right", BLOOD_RED)
    add_top_bottom_bars(slide, DARK_BURGUNDY, Inches(0.06))

    add_centered_heading(slide, "The Traitors", y_pos=Inches(0.4), color=BURGUNDY,
                          size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(1.15), BURGUNDY, width=Inches(4))

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "They are already among you.",
             "size": Pt(32), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(28)},
            {"text": "6 players have been chosen to deceive.",
             "size": Pt(28), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(24)},
            {"text": "They will participate in every challenge.",
             "size": Pt(26), "color": LIGHT_GRAY, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(10)},
            {"text": "They will debate alongside you at the Round Table.",
             "size": Pt(26), "color": LIGHT_GRAY, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(10)},
            {"text": "They will vote at every banishment.",
             "size": Pt(26), "color": LIGHT_GRAY, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(10)},
            {"text": "They will look you in the eye and lie.",
             "size": Pt(26), "color": LIGHT_GRAY, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(30)},
            {"text": "And when night falls...",
             "size": Pt(30), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(10)},
            {"text": "someone will not return.",
             "size": Pt(30), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "Keep this vague and menacing. Do NOT reveal any traitor mechanics.\n\n"
        "'Six of you have been chosen. You know who you are. You will walk among "
        "the Faithful. You will participate in every challenge, sit at every debate, "
        "cast your vote at every banishment -- all while hiding what you truly are. "
        "And when night falls... you will decide who does not return.\n\n"
        "To the Faithful: they are already watching you. To the Traitors: "
        "the game has begun.'\n\n"
        "Do NOT elaborate on traitor powers, night mechanics, or how traitors "
        "communicate. Keep the mystery alive."
    ))


def build_slide_07_houses(prs):
    """Slide 7: The Houses - House A and House B explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NEAR_BLACK)
    add_side_accent(slide, "left", BURGUNDY)
    add_side_accent(slide, "right", BURGUNDY)
    add_top_bottom_bars(slide, DARK_BURGUNDY)

    add_centered_heading(slide, "The Houses", y_pos=Inches(0.3), color=GOLD)
    add_decorative_bar(slide, Inches(1.05), GOLD, width=Inches(4))

    # Intro text
    add_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(10)) / 2,
        top=Inches(1.3),
        width=Inches(10),
        height=Inches(0.8),
        text="You will be divided into two rival houses of 17.",
        font_size=Pt(28), font_color=CREAM, alignment=PP_ALIGN.CENTER
    )

    # House A box
    box_width = Inches(4.5)
    box_height = Inches(3.5)
    gap = Inches(1)
    total_w = box_width * 2 + gap
    start_x = (SLIDE_WIDTH - total_w) / 2

    house_a = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        start_x, Inches(2.4), box_width, box_height
    )
    house_a.fill.solid()
    house_a.fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x0A)
    house_a.line.color.rgb = BURGUNDY
    house_a.line.width = Pt(2)

    add_multi_text_box(
        slide,
        left=start_x + Inches(0.3),
        top=Inches(2.6),
        width=box_width - Inches(0.6),
        height=box_height - Inches(0.4),
        lines=[
            {"text": "HOUSE A", "size": Pt(36), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(16)},
            {"text": "17 Players", "size": Pt(26), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "3 Traitors hidden within", "size": Pt(22), "color": DIM_GRAY,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(16)},
            {"text": "Traitors in House A know\nonly each other.",
             "size": Pt(20), "color": LIGHT_GRAY,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    # House B box
    house_b_x = start_x + box_width + gap
    house_b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        house_b_x, Inches(2.4), box_width, box_height
    )
    house_b.fill.solid()
    house_b.fill.fore_color.rgb = RGBColor(0x0A, 0x14, 0x0E)
    house_b.line.color.rgb = DARK_GREEN
    house_b.line.width = Pt(2)

    add_multi_text_box(
        slide,
        left=house_b_x + Inches(0.3),
        top=Inches(2.6),
        width=box_width - Inches(0.6),
        height=box_height - Inches(0.4),
        lines=[
            {"text": "HOUSE B", "size": Pt(36), "color": MEDIUM_GREEN, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(16)},
            {"text": "17 Players", "size": Pt(26), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "3 Traitors hidden within", "size": Pt(22), "color": DIM_GRAY,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(16)},
            {"text": "Traitors in House B know\nonly each other.",
             "size": Pt(20), "color": LIGHT_GRAY,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    # Bottom note
    add_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(10)) / 2,
        top=Inches(6.2),
        width=Inches(10),
        height=Inches(0.8),
        text="The houses will merge when 22 players remain.",
        font_size=Pt(26), font_color=GOLD, bold=True, alignment=PP_ALIGN.CENTER
    )

    add_speaker_notes(slide, (
        "'You will be split into two houses -- House A and House B -- each with 17 "
        "players. Each house contains 3 Traitors. During Phase 1, you will only "
        "interact with your own house. The Traitors in House A only know the other "
        "Traitors in House A. Same for House B. They do not know who the Traitors "
        "are in the other house.\n\n"
        "When we are down to 22 players, the houses will merge into one. "
        "At that point... things get very interesting.'\n\n"
        "NOTE: House assignments were included in the envelopes."
    ))


def build_slide_08_phase1(prs):
    """Slide 8: Phase 1: The Houses - how house rounds work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    standard_slide_setup(slide, "Phase 1: The Houses")

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    # Round structure summary
    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.5),
        width=content_width,
        height=Inches(1.2),
        lines=[
            {"text": "3 rotations per house  |  34 players down to 22",
             "size": Pt(26), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(14)},
            {"text": "Each rotation follows the same deadly cycle:",
             "size": Pt(24), "color": CREAM, "alignment": PP_ALIGN.CENTER},
        ]
    )

    # Four phase boxes
    steps = [
        ("THE CHALLENGE", "5 MIN", BURGUNDY, "Compete.\nProve yourself."),
        ("THE DEBATE", "10 MIN", DARK_GREEN, "Accuse. Defend.\nPersuade."),
        ("THE VOTE", "5 MIN", GOLD, "Banish one\nplayer."),
        ("THE NIGHT", "2 MIN", RGBColor(0x44, 0x00, 0x44), "Someone is\nmurdered."),
    ]

    box_w = Inches(2.5)
    box_h = Inches(3.0)
    gap = Inches(0.3)
    total_w = box_w * 4 + gap * 3
    start_x = (SLIDE_WIDTH - total_w) / 2

    for i, (title, time_str, color, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        y = Inches(3.1)

        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x18)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Colored title bar inside box
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Title text on bar
        title_color = WHITE if color != GOLD else BLACK
        add_text_box(
            slide, x + Inches(0.1), y + Inches(0.05),
            box_w - Inches(0.2), Inches(0.45),
            title, font_size=Pt(16), font_color=title_color,
            bold=True, alignment=PP_ALIGN.CENTER
        )

        # Time and description
        add_multi_text_box(
            slide,
            left=x + Inches(0.2),
            top=y + Inches(0.7),
            width=box_w - Inches(0.4),
            height=Inches(2.0),
            lines=[
                {"text": time_str, "size": Pt(28), "color": color, "bold": True,
                 "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(12)},
                {"text": desc, "size": Pt(18), "color": LIGHT_GRAY,
                 "alignment": PP_ALIGN.CENTER},
            ]
        )

    # Bottom summary
    add_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(10)) / 2,
        top=Inches(6.3),
        width=Inches(10),
        height=Inches(0.7),
        text="Each cycle removes 2 players per house (1 banished + 1 murdered)",
        font_size=Pt(22), font_color=DIM_GRAY, alignment=PP_ALIGN.CENTER
    )

    add_speaker_notes(slide, (
        "'Phase 1 is the House phase. Each house runs through 3 complete rotations. "
        "Each rotation has four parts: a 5-minute Challenge, a 10-minute Debate at "
        "the Round Table, a 5-minute Vote to banish someone, and a 2-minute Night "
        "phase where someone is murdered.\n\n"
        "That means every single rotation, two people leave the game from each house "
        "-- one voted out by the group, one eliminated in the dark. After 3 "
        "rotations, we go from 34 down to 22. At 22, the houses merge.\n\n"
        "Houses alternate. House A goes first each rotation, then House B.'"
    ))


def build_slide_09_challenge(prs):
    """Slide 9: The Challenge - how challenges work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    standard_slide_setup(slide, "The Challenge")

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "Each round begins with a challenge.",
             "size": Pt(30), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": diamond_bullet("Challenges test knowledge, speed, or teamwork"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(12)},
            {"text": diamond_bullet("Each challenge has a winner or winning team"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(12)},
            {"text": diamond_bullet("Winning may grant protection or special powers"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(12)},
            {"text": diamond_bullet("Challenge results feed into the debate"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(28)},
            {"text": "Watch carefully.",
             "size": Pt(30), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(10)},
            {"text": "Who tried to win? Who held back? Who sabotaged?",
             "size": Pt(24), "color": DIM_GRAY,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "'Every round starts with a challenge. These are quick -- 5 minutes or less. "
        "They might be trivia, physical tasks, or team puzzles. Pay attention to who "
        "is trying to win and who might be quietly holding back. Challenge winners "
        "may earn special advantages -- sometimes protection, sometimes a power that "
        "can change the game.\n\n"
        "I will explain each challenge before we play it. The rules will be simple. "
        "But remember: everything in this game is information. How someone plays a "
        "challenge tells you something about who they are.'"
    ))


def build_slide_10_round_table(prs):
    """Slide 10: The Round Table - debate rules, 30-second speaking limit."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NEAR_BLACK)

    # Dark green table feel
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0A, 0x14, 0x0D)
    overlay.line.fill.background()

    add_side_accent(slide, "left", DARK_GREEN)
    add_side_accent(slide, "right", DARK_GREEN)
    add_top_bottom_bars(slide, DARK_GREEN, Inches(0.06))

    add_centered_heading(slide, "The Round Table", y_pos=Inches(0.35),
                          color=GOLD, size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(1.1), GOLD, width=Inches(5))

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.5),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "This is where the game is won and lost.",
             "size": Pt(28), "color": CREAM, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(20)},
            {"text": "THE RULES OF DEBATE", "size": Pt(26), "color": GOLD,
             "bold": True, "spacing_after": Pt(16)},
            {"text": diamond_bullet("You have 30 SECONDS when called upon to speak"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("The Host controls who speaks and when"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("You may accuse, defend, question, or deflect"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("NO WHISPERING -- all conversation is public"),
             "size": BODY_SIZE, "color": BURGUNDY, "bold": True,
             "spacing_after": Pt(10)},
            {"text": diamond_bullet("No phones during the debate"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(20)},
            {"text": "Speak with conviction. Listen for lies.",
             "size": Pt(28), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "10 minutes. Make every second count.",
             "size": Pt(24), "color": DIM_GRAY,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "'After each challenge, we gather at the Round Table for a 10-minute debate. "
        "This is the heart of the game. When I call on you, you have exactly 30 "
        "seconds to speak. You can accuse someone, defend yourself, ask a question, "
        "or make an argument. When your 30 seconds are up, you stop.\n\n"
        "CRITICAL RULES: All conversation must be public. No whispering. No side "
        "conversations. No phones. If I catch whispering, I will call it out. "
        "Everything said at this table is heard by everyone.\n\n"
        "This is where Traitors will try to blend in -- and where Faithful players "
        "will try to catch them. Listen carefully. Watch body language. "
        "The truth is in the details.'\n\n"
        "HOST TIP: Use a visible timer. Call on players who are too quiet. "
        "Keep energy high. Cut people off at 30 seconds firmly but fairly."
    ))


def build_slide_11_banishment(prs):
    """Slide 11: The Banishment - how voting works."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    standard_slide_setup(slide, "The Banishment")

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "After the debate, you vote.",
             "size": Pt(30), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(20)},
            {"text": "HOW IT WORKS", "size": Pt(26), "color": GOLD, "bold": True,
             "spacing_after": Pt(14)},
            {"text": diamond_bullet("Everyone writes one name on their voting slip"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Votes are collected and read aloud by the Host"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("The player with the most votes is BANISHED"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("In case of a tie, the Host calls a revote "
                                     "between tied players"),
             "size": SMALL_BODY_SIZE, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
            {"text": diamond_bullet("If still tied, the Host breaks the tie"),
             "size": SMALL_BODY_SIZE, "color": LIGHT_GRAY, "spacing_after": Pt(20)},
            {"text": "The banished player reveals their role.",
             "size": Pt(28), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "Faithful or Traitor -- the truth comes out.",
             "size": Pt(24), "color": DIM_GRAY,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "'After the debate, it is time to vote. Each of you writes one name on your "
        "voting slip -- the person you want banished. No discussion during voting. "
        "I will collect them, read them aloud one by one, and the player with the "
        "most votes is banished from the game.\n\n"
        "If there is a tie, we do a quick revote between just the tied players. "
        "If still tied, I break the tie as Host.\n\n"
        "When you are banished, you reveal your role card. You show the room whether "
        "you were Faithful or a Traitor. This is a critical moment -- it gives "
        "everyone information. Did the group make the right call, or did the "
        "Traitors just convince you to eliminate one of your own?\n\n"
        "Choose carefully. Every wrong vote is a win for the Traitors.'"
    ))


def build_slide_12_night(prs):
    """Slide 12: The Night - mysterious, no mechanical details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)

    # Very dark purple/black overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0D, 0x06, 0x12)
    overlay.line.fill.background()

    add_side_accent(slide, "left", RGBColor(0x33, 0x00, 0x33))
    add_side_accent(slide, "right", RGBColor(0x33, 0x00, 0x33))
    add_top_bottom_bars(slide, RGBColor(0x33, 0x00, 0x33), Inches(0.06))

    add_centered_heading(slide, "The Night", y_pos=Inches(0.5),
                          color=PALE_GOLD, size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(1.25), RGBColor(0x66, 0x33, 0x66),
                        width=Inches(3))

    add_multi_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(9)) / 2,
        top=Inches(1.8),
        width=Inches(9),
        height=Inches(5),
        lines=[
            {"text": "When darkness falls, the Traitors go to work.",
             "size": Pt(30), "color": CREAM, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(28)},
            {"text": diamond_bullet("All players close their eyes or bow their heads"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(12)},
            {"text": diamond_bullet("The Host will conduct the night phase"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(12)},
            {"text": diamond_bullet("When morning comes, one player will be gone"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(12)},
            {"text": diamond_bullet("The murdered player reveals their role"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(28)},
            {"text": "You will not know how it happened.",
             "size": Pt(28), "color": RGBColor(0x99, 0x66, 0x99),
             "bold": True, "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(10)},
            {"text": "You will only know they are gone.",
             "size": Pt(28), "color": RGBColor(0x99, 0x66, 0x99),
             "bold": True, "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "Keep this mysterious. Do NOT explain how traitors choose their victim.\n\n"
        "'After the banishment, night falls. I will ask everyone to close their eyes "
        "or bow their heads. The room must be silent. I will conduct the night phase. "
        "You do not need to know what happens in the dark -- only that when morning "
        "comes, someone will not be with us anymore. The murdered player will be "
        "revealed and their role shown.\n\n"
        "Note: Phones may ONLY be used during the night phase if needed. "
        "Put them away when morning comes.'\n\n"
        "HOST MECHANICS (do not say aloud): Tap traitors on shoulder to open eyes. "
        "They silently point to agree on a victim. You confirm with a nod. "
        "They close eyes. You tap the victim. Morning."
    ))


def build_slide_13_merge(prs):
    """Slide 13: Phase 2: The Merge - when houses combine."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NEAR_BLACK)

    # Dramatic dual-tone overlay
    left_half = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH // 2, SLIDE_HEIGHT
    )
    left_half.fill.solid()
    left_half.fill.fore_color.rgb = RGBColor(0x14, 0x08, 0x08)
    left_half.line.fill.background()

    right_half = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, SLIDE_WIDTH // 2, 0, SLIDE_WIDTH // 2, SLIDE_HEIGHT
    )
    right_half.fill.solid()
    right_half.fill.fore_color.rgb = RGBColor(0x08, 0x14, 0x0A)
    right_half.line.fill.background()

    # Center overlay to blend the halves
    center_blend = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), 0, SLIDE_WIDTH - Inches(4), SLIDE_HEIGHT
    )
    center_blend.fill.solid()
    center_blend.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x10)
    center_blend.line.fill.background()

    add_side_accent(slide, "left", BURGUNDY)
    add_side_accent(slide, "right", DARK_GREEN)
    add_top_bottom_bars(slide, GOLD, Inches(0.06))

    add_centered_heading(slide, "Phase 2: The Merge", y_pos=Inches(0.4),
                          color=GOLD, size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(1.15), GOLD, width=Inches(5))

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "At 22 players, the houses become one.",
             "size": Pt(32), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": diamond_bullet("House A and House B are dissolved"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("All remaining players form a single group"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("All surviving Traitors now know each other"),
             "size": BODY_SIZE, "color": BURGUNDY, "bold": True,
             "spacing_after": Pt(20)},
            {"text": "THE MERGE ROUND", "size": Pt(26), "color": GOLD, "bold": True,
             "spacing_after": Pt(14)},
            {"text": diamond_bullet("A special Merge Challenge for a powerful reward"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("The Chalice Ceremony -- a toast to new alliances"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("A DOUBLE banishment -- two players voted out"),
             "size": BODY_SIZE, "color": BLOOD_RED, "bold": True,
             "spacing_after": Pt(10)},
            {"text": diamond_bullet("Then the night..."),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
        ]
    )

    add_speaker_notes(slide, (
        "'When we reach 22 players, the two houses merge into one. House A and "
        "House B are dissolved. You are all together now. And here is the key "
        "moment: ALL surviving Traitors now know each other. Before the merge, "
        "Traitors only knew their own housemates. Now the full network is revealed "
        "-- to them.\n\n"
        "The merge round is special. There will be a Merge Challenge with a powerful "
        "reward -- the Dagger. There will be the Chalice Ceremony -- a ritual toast. "
        "And there will be a DOUBLE banishment -- two players will be voted out in "
        "the same round. Then night falls as usual.\n\n"
        "From the merge onward, the game gets faster and more intense. "
        "22 players to the endgame.'\n\n"
        "NOTE: If fewer than 3 Traitors remain at merge, secretly recruit 1. "
        "Do this privately before the merge announcement."
    ))


def build_slide_14_dagger(prs):
    """Slide 14: The Dagger - what it does, how it is earned."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)

    # Warm golden overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x12, 0x0C, 0x00)
    overlay.line.fill.background()

    add_side_accent(slide, "left", GOLD)
    add_side_accent(slide, "right", GOLD)
    add_top_bottom_bars(slide, RGBColor(0x8B, 0x75, 0x00), Inches(0.06))

    # Dagger icon (diamond shape)
    dagger_icon = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        (SLIDE_WIDTH - Inches(0.8)) / 2,
        Inches(0.4),
        Inches(0.8),
        Inches(1.0)
    )
    dagger_icon.fill.solid()
    dagger_icon.fill.fore_color.rgb = GOLD
    dagger_icon.line.fill.background()

    add_centered_heading(slide, "The Dagger", y_pos=Inches(1.5),
                          color=GOLD, size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(2.2), GOLD, width=Inches(4))

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(2.5),
        width=content_width,
        height=Inches(4.5),
        lines=[
            {"text": "Won by the victor of the Merge Challenge.",
             "size": Pt(28), "color": CREAM, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": "THE POWER", "size": Pt(26), "color": GOLD, "bold": True,
             "spacing_after": Pt(14)},
            {"text": diamond_bullet("Cast TWO votes in a single banishment"),
             "size": Pt(30), "color": PALE_GOLD, "bold": True,
             "spacing_after": Pt(20)},
            {"text": "THE RULES", "size": Pt(26), "color": GOLD, "bold": True,
             "spacing_after": Pt(14)},
            {"text": diamond_bullet("May be used this round or the next"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(8)},
            {"text": diamond_bullet("Must declare BEFORE votes are revealed"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(8)},
            {"text": diamond_bullet("Both votes must target the same player"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(8)},
            {"text": diamond_bullet("Expires after 2 rounds -- use it or lose it"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(8)},
            {"text": diamond_bullet("Cannot be used at the Final 6"),
             "size": SMALL_BODY_SIZE, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
            {"text": diamond_bullet("If the holder is banished, the Dagger vanishes"),
             "size": SMALL_BODY_SIZE, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
        ]
    )

    add_speaker_notes(slide, (
        "'The Merge Challenge awards a powerful weapon: The Dagger. Whoever wins the "
        "challenge earns the right to cast TWO votes in a single banishment instead "
        "of one. Both votes must go to the same person.\n\n"
        "You must announce that you are using the Dagger BEFORE the votes are "
        "revealed. You can use it the round you win it or save it for the next "
        "round -- but after two rounds, it expires. And if you get banished while "
        "holding it, the Dagger is gone.\n\n"
        "This is political power. It can swing a close vote. It makes you a target "
        "AND an ally. Choose wisely when to play it.'\n\n"
        "HOST NOTE: The Dagger cannot be used at the Final 6. "
        "Announce dagger use publicly before reading votes."
    ))


def build_slide_15_chalice(prs):
    """Slide 15: The Chalice Ceremony - atmospheric, NO secret mechanics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)

    # Rich warm overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x14, 0x0A, 0x02)
    overlay.line.fill.background()

    add_side_accent(slide, "left", BURGUNDY)
    add_side_accent(slide, "right", BURGUNDY)
    add_top_bottom_bars(slide, DARK_BURGUNDY, Inches(0.06))

    # Chalice icon (oval with gold border)
    chalice_icon = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        (SLIDE_WIDTH - Inches(1.0)) / 2,
        Inches(0.5),
        Inches(1.0),
        Inches(1.0)
    )
    chalice_icon.fill.solid()
    chalice_icon.fill.fore_color.rgb = DARK_BURGUNDY
    chalice_icon.line.color.rgb = GOLD
    chalice_icon.line.width = Pt(2)

    add_centered_heading(slide, "The Chalice Ceremony", y_pos=Inches(1.7),
                          color=GOLD, size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(2.45), GOLD, width=Inches(5))

    add_multi_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(9)) / 2,
        top=Inches(2.8),
        width=Inches(9),
        height=Inches(4),
        lines=[
            {"text": "A toast to new alliances.",
             "size": Pt(32), "color": PALE_GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(28)},
            {"text": "When the houses merge, we gather to mark the occasion.",
             "size": Pt(26), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(16)},
            {"text": "Each player will raise a chalice.",
             "size": Pt(26), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(16)},
            {"text": "You will mingle. You will toast. You will drink.",
             "size": Pt(26), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(28)},
            {"text": "Enjoy the ceremony.",
             "size": Pt(30), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(12)},
            {"text": "But be careful who you drink with.",
             "size": Pt(28), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "THIS SLIDE IS DELIBERATELY VAGUE. Do NOT explain the Chalice Mark mechanic. "
        "The secret mechanic is that one Traitor can 'mark' a player during the "
        "ceremony by interlocking arms and saying 'Cheers' -- that player is killed "
        "at night instead of the normal murder. But players must NOT know this.\n\n"
        "Say: 'To celebrate the merge, we hold the Chalice Ceremony. Each of you "
        "will receive a chalice. You will have 3 to 4 minutes to mingle, chat, and "
        "toast to new alliances. It is a moment of celebration... and perhaps "
        "something more. Enjoy it. But be mindful of who you share your toast with.'\n\n"
        "The vagueness IS the point. Players will wonder later why you mentioned "
        "'be careful who you drink with.' That seed of doubt is intentional.\n\n"
        "HOST: Watch for the traitor arm-interlock. You must witness it for it to count."
    ))


def build_slide_16_endgame(prs):
    """Slide 16: The Endgame - from 11 to 6, banishments only."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    standard_slide_setup(slide, "The Endgame")

    content_left = Inches(1.5)
    content_width = SLIDE_WIDTH - Inches(3)

    add_multi_text_box(
        slide,
        left=content_left,
        top=Inches(1.6),
        width=content_width,
        height=Inches(5.5),
        lines=[
            {"text": "At 11 players, the murders stop.",
             "size": Pt(32), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": "From this point forward, only banishments remain.",
             "size": Pt(26), "color": CREAM, "alignment": PP_ALIGN.CENTER,
             "spacing_after": Pt(24)},
            {"text": diamond_bullet("No more night kills"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Each round: debate, then vote"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("One player banished per round"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(10)},
            {"text": diamond_bullet("Banished players still reveal their role"),
             "size": BODY_SIZE, "color": CREAM, "spacing_after": Pt(24)},
            {"text": "11 ... 10 ... 9 ... 8 ... 7 ... 6",
             "size": Pt(36), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(16)},
            {"text": "At 6, the Final Circle is called.",
             "size": Pt(28), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "'When we reach 11 players, the game changes. No more murders. No more "
        "night phase. From 11 down to 6, the only way someone leaves is by "
        "banishment. That means every single vote counts even more. There is "
        "nowhere to hide.\n\n"
        "Each round is a debate followed by a vote. One person is banished. "
        "They reveal their role. And the circle gets smaller. 11, 10, 9, 8, 7... "
        "until we reach 6.\n\n"
        "At 6 players, the Final Circle is called. This is the moment everything "
        "has been building toward.'"
    ))


def build_slide_17_final_circle(prs):
    """Slide 17: The Final Circle - the final vote to end or continue."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)

    # Dramatic dark red overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x1A, 0x05, 0x05)
    overlay.line.fill.background()

    # Gold border frame
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.4),
        SLIDE_WIDTH - Inches(1.6), SLIDE_HEIGHT - Inches(0.8)
    )
    frame.fill.background()
    frame.line.color.rgb = GOLD
    frame.line.width = Pt(2)

    add_side_accent(slide, "left", GOLD)
    add_side_accent(slide, "right", GOLD)

    add_centered_heading(slide, "The Final Circle", y_pos=Inches(0.6),
                          color=GOLD, size=LARGE_TITLE_SIZE)
    add_decorative_bar(slide, Inches(1.5), GOLD, width=Inches(5))

    add_multi_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(9)) / 2,
        top=Inches(1.9),
        width=Inches(9),
        height=Inches(5),
        lines=[
            {"text": "6 players remain. The moment of truth.",
             "size": Pt(30), "color": CREAM, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": "The group must make a choice:",
             "size": Pt(26), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": "\"END THE GAME\"",
             "size": Pt(36), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(10)},
            {"text": "If the group votes to END, the game is over.",
             "size": Pt(24), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "All remaining Traitors are revealed.",
             "size": Pt(24), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(20)},
            {"text": "\"CONTINUE\"",
             "size": Pt(36), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(10)},
            {"text": "If the group votes to CONTINUE, another banishment occurs.",
             "size": Pt(24), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "The circle shrinks further. Another vote at 5, at 4...",
             "size": Pt(24), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(24)},
            {"text": "If any Traitors remain when you END: they win.",
             "size": Pt(28), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
            {"text": "If all Traitors are gone when you END: the Faithful win.",
             "size": Pt(28), "color": GOLD, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "'At 6 players, I will offer the group a choice. You can vote to END THE "
        "GAME -- right now -- or you can vote to CONTINUE. If you end the game, "
        "all remaining Traitors reveal themselves. If even ONE Traitor is still "
        "standing, the Traitors win. If all Traitors have been banished, the "
        "Faithful win.\n\n"
        "But if you are not confident -- if you think a Traitor is still hiding -- "
        "you can vote to CONTINUE. This triggers another banishment round. You "
        "shrink to 5, then potentially 4, and each time you can vote to end or "
        "continue.\n\n"
        "The tension here is enormous. End too early and a Traitor slips through. "
        "Continue too long and you might banish your own allies.\n\n"
        "This is the final test of trust.'\n\n"
        "HOST: A simple majority vote decides END vs CONTINUE. "
        "Ties default to CONTINUE."
    ))


def build_slide_18_rules(prs):
    """Slide 18: The Rules - critical gameplay rules as formatted blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NEAR_BLACK)

    add_side_accent(slide, "left", BURGUNDY)
    add_side_accent(slide, "right", BURGUNDY)
    add_top_bottom_bars(slide, BURGUNDY, Inches(0.08))

    add_centered_heading(slide, "The Rules", y_pos=Inches(0.3), color=GOLD,
                          size=HEADING_SIZE)
    add_decorative_bar(slide, Inches(1.05), GOLD, width=Inches(3))

    # Rules as distinct visual blocks
    rules = [
        ("30-SECOND SPEAKING LIMIT",
         "When called upon, you have 30 seconds. No more.",
         BURGUNDY),
        ("NO WHISPERING",
         "All conversation is public. Side chats are forbidden.",
         BURGUNDY),
        ("NO PHONES",
         "Phones away during all active phases. Only permitted during Night.",
         BURGUNDY),
        ("ELIMINATED PLAYERS: SILENCE",
         "Once banished or murdered, you may not speak, gesture, or influence.",
         BURGUNDY),
        ("HOST DECISIONS ARE FINAL",
         "The Host's word is law. No appeals. No arguments.",
         GOLD),
    ]

    start_y = Inches(1.4)
    rule_height = Inches(1.05)
    rule_width = Inches(10)
    rule_x = (SLIDE_WIDTH - rule_width) / 2

    for i, (title, desc, accent) in enumerate(rules):
        y = start_y + i * (rule_height + Inches(0.08))

        # Rule background
        rule_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, rule_x, y, rule_width, rule_height
        )
        rule_bg.fill.solid()
        rule_bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
        rule_bg.line.fill.background()

        # Left accent bar on block
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, rule_x, y, Inches(0.12), rule_height
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # Rule title
        add_text_box(
            slide, rule_x + Inches(0.4), y + Inches(0.08),
            rule_width - Inches(0.8), Inches(0.45),
            title, font_size=Pt(24), font_color=accent, bold=True
        )

        # Rule description
        add_text_box(
            slide, rule_x + Inches(0.4), y + Inches(0.52),
            rule_width - Inches(0.8), Inches(0.45),
            desc, font_size=Pt(20), font_color=LIGHT_GRAY
        )

    add_speaker_notes(slide, (
        "Read each rule clearly and firmly. These are non-negotiable.\n\n"
        "'Before we begin, the rules. These are absolute.\n\n"
        "ONE: You have 30 seconds to speak when called upon. I will cut you off.\n\n"
        "TWO: No whispering. All conversation happens publicly. If I see or hear "
        "side conversations, I will call you out.\n\n"
        "THREE: Phones away. No texting, no Googling, no messaging. The only time "
        "phones are permitted is during the Night phase.\n\n"
        "FOUR: If you are eliminated -- banished or murdered -- you are silent. "
        "No speaking, no gesturing, no mouthing words, no facial expressions "
        "meant to influence. You are a ghost.\n\n"
        "FIVE: My decisions as Host are final. If there is a dispute, I rule. "
        "No appeals.\n\n"
        "These rules exist to keep the game fair and dramatic. Respect them.'"
    ))


def build_slide_19_begin(prs):
    """Slide 19: Let the Game Begin - dramatic closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)

    # Deep atmospheric layers
    layer1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    layer1.fill.solid()
    layer1.fill.fore_color.rgb = RGBColor(0x0A, 0x04, 0x04)
    layer1.line.fill.background()

    # Inner dramatic frame with gold border
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(0.8),
        SLIDE_WIDTH - Inches(3), SLIDE_HEIGHT - Inches(1.6)
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x0E, 0x06, 0x06)
    frame.line.color.rgb = GOLD
    frame.line.width = Pt(1.5)

    add_top_bottom_bars(slide, BURGUNDY, Inches(0.06))

    # Top decorative bar
    add_decorative_bar(slide, Inches(2.3), GOLD, width=Inches(5), height=Inches(0.03))

    # Main text
    add_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(10)) / 2,
        top=Inches(2.5),
        width=Inches(10),
        height=Inches(1.5),
        text="Let the Game Begin",
        font_size=LARGE_TITLE_SIZE,
        font_color=GOLD,
        bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Bottom decorative bar
    add_decorative_bar(slide, Inches(4.1), GOLD, width=Inches(5), height=Inches(0.03))

    # Atmospheric closing lines
    add_multi_text_box(
        slide,
        left=(SLIDE_WIDTH - Inches(8)) / 2,
        top=Inches(4.4),
        width=Inches(8),
        height=Inches(2.5),
        lines=[
            {"text": "The castle doors are sealed.",
             "size": Pt(26), "color": CREAM,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(14)},
            {"text": "The Traitors are among you.",
             "size": Pt(26), "color": BURGUNDY, "bold": True,
             "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(14)},
            {"text": "Trust no one.",
             "size": Pt(32), "color": BLOOD_RED, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ]
    )

    add_speaker_notes(slide, (
        "This is the final slide before gameplay begins. Build the tension.\n\n"
        "'The rules have been set. Your roles have been assigned. The houses have "
        "been formed. There is nothing left to explain.\n\n"
        "[Pause]\n\n"
        "The castle doors are sealed. The Traitors are among you. "
        "From this moment on, every word, every glance, every silence... "
        "means something.\n\n"
        "[Long pause]\n\n"
        "Trust no one. Let the game begin.'\n\n"
        "HOST: Immediately transition to Phase 1. Call House A to their area. "
        "Begin the first challenge. Keep the energy high. The game is ON."
    ))


# ==============================================================================
# MAIN
# ==============================================================================

def create_presentation():
    """Build and save the complete Traitors presentation."""
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Building THE TRAITORS presentation...")
    print("=" * 50)

    builders = [
        ("Slide 01: Title", build_slide_01_title),
        ("Slide 02: Welcome to the Castle", build_slide_02_welcome),
        ("Slide 03: The Game", build_slide_03_the_game),
        ("Slide 04: Your Role", build_slide_04_your_role),
        ("Slide 05: The Faithful", build_slide_05_faithful),
        ("Slide 06: The Traitors", build_slide_06_traitors),
        ("Slide 07: The Houses", build_slide_07_houses),
        ("Slide 08: Phase 1 - The Houses", build_slide_08_phase1),
        ("Slide 09: The Challenge", build_slide_09_challenge),
        ("Slide 10: The Round Table", build_slide_10_round_table),
        ("Slide 11: The Banishment", build_slide_11_banishment),
        ("Slide 12: The Night", build_slide_12_night),
        ("Slide 13: Phase 2 - The Merge", build_slide_13_merge),
        ("Slide 14: The Dagger", build_slide_14_dagger),
        ("Slide 15: The Chalice Ceremony", build_slide_15_chalice),
        ("Slide 16: The Endgame", build_slide_16_endgame),
        ("Slide 17: The Final Circle", build_slide_17_final_circle),
        ("Slide 18: The Rules", build_slide_18_rules),
        ("Slide 19: Let the Game Begin", build_slide_19_begin),
    ]

    for name, builder in builders:
        print(f"  Building {name}...")
        builder(prs)

    output_dir = "/Users/johann.kerr/Documents/traitors"
    output_path = os.path.join(output_dir, "traitors_rules.pptx")
    prs.save(output_path)

    print("=" * 50)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nSlide summary:")
    for i, (name, _) in enumerate(builders, 1):
        print(f"  {i:2d}. {name.split(': ', 1)[1]}")
    print("\nDone! Open in PowerPoint or Keynote to review.")


if __name__ == "__main__":
    create_presentation()
